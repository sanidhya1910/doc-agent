"""Turn anything the user drops in into a :class:`~docagent.state.Document`.

Handles images (including multi-frame TIFF and HEIC), PDFs, Office files and
ZIP archives.  Two decisions here matter for cost:

* PDFs are rendered with ``pypdfium2``, which bundles its own PDFium.  That
  avoids a ``packages.txt`` poppler install and keeps Space builds fast.
* When a PDF or Office file already carries text, it is used as-is and the
  page is marked ``has_text_layer``.  Those pages never reach the GPU, which
  is the single largest saving against the ZeroGPU daily quota.
"""

from __future__ import annotations

import logging
import os
import zipfile
from pathlib import Path
from typing import Any, Sequence

from PIL import Image, ImageSequence

from .state import Block, Document, Page, Table
from .tables import table_to_markdown

log = logging.getLogger("docagent.ingest")

DEFAULT_DPI = 200
#: Fewer characters than this on a PDF page means the "text layer" is really
#: just a header or a stray annotation over a scan, so OCR is still needed.
TEXT_LAYER_MIN_CHARS = 40

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff", ".ppm"}
HEIC_SUFFIXES = {".heic", ".heif"}
PDF_SUFFIXES = {".pdf"}
DOCX_SUFFIXES = {".docx"}
PPTX_SUFFIXES = {".pptx"}
XLSX_SUFFIXES = {".xlsx", ".xlsm"}
TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".csv"}
ZIP_SUFFIXES = {".zip"}

SUPPORTED_SUFFIXES = (
    IMAGE_SUFFIXES | HEIC_SUFFIXES | PDF_SUFFIXES | DOCX_SUFFIXES
    | PPTX_SUFFIXES | XLSX_SUFFIXES | TEXT_SUFFIXES | ZIP_SUFFIXES
)


def _register_heif() -> bool:
    try:
        import pillow_heif

        pillow_heif.register_heif_opener()
        return True
    except Exception:
        return False


# ---------------------------------------------------------------------------
# per-format loaders, each yielding Page objects with index left at 0
# ---------------------------------------------------------------------------

def _pages_from_image(path: Path) -> list[Page]:
    if path.suffix.lower() in HEIC_SUFFIXES and not _register_heif():
        raise ValueError("HEIC support needs pillow-heif; install it or convert the file first")

    out: list[Page] = []
    with Image.open(path) as img:
        # Multi-frame TIFF and animated GIF become one page per frame.
        for frame in ImageSequence.Iterator(img):
            out.append(Page(index=0, image=frame.convert("RGB"), source=path.name))
    return out


def _pages_from_pdf(path: Path, dpi: int) -> list[Page]:
    import pypdfium2 as pdfium

    out: list[Page] = []
    pdf = pdfium.PdfDocument(str(path))
    try:
        scale = dpi / 72.0
        for i in range(len(pdf)):
            page = pdf[i]
            text = ""
            try:
                textpage = page.get_textpage()
                text = (textpage.get_text_range() or "").strip()
                textpage.close()
            except Exception as exc:
                log.debug("no text layer on %s p%d: %s", path.name, i + 1, exc)

            image = page.render(scale=scale).to_pil().convert("RGB")
            has_text = len(text) >= TEXT_LAYER_MIN_CHARS
            out.append(
                Page(
                    index=0,
                    image=image,
                    source=path.name,
                    markdown=text if has_text else "",
                    blocks=[Block(text=line) for line in text.splitlines() if line.strip()]
                    if has_text
                    else [],
                    has_text_layer=has_text,
                    backend="pdf_text_layer" if has_text else "",
                )
            )
            page.close()
    finally:
        pdf.close()
    return out


def _with_tables(lines: list[str], tables: list[Table]) -> str:
    """Page text with its parsed tables rendered back in as Markdown.

    Office formats hand us tables as structure, not prose, so without this a
    slide or document whose content is a table exports empty text.
    """
    parts = [ln for ln in lines if ln.strip()]
    parts.extend(table_to_markdown(t) for t in tables)
    return "\n\n".join(parts)


def _table_from_rows(rows: list[list[str]]) -> Table | None:
    rows = [r for r in rows if any(c.strip() for c in r)]
    if len(rows) < 2:
        return None
    return Table(header=rows[0], rows=rows[1:])


def _pages_from_docx(path: Path) -> list[Page]:
    import docx

    document = docx.Document(str(path))
    lines: list[str] = []
    for para in document.paragraphs:
        if not para.text.strip():
            continue
        level = 0
        if para.style is not None and (para.style.name or "").startswith("Heading"):
            tail = (para.style.name or "").split()[-1]
            level = int(tail) if tail.isdigit() else 1
        lines.append(("#" * level + " " if level else "") + para.text.strip())

    tables: list[Table] = []
    for tbl in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
        parsed = _table_from_rows(rows)
        if parsed:
            tables.append(parsed)

    return [
        Page(
            index=0,
            image=None,
            source=path.name,
            markdown=_with_tables(lines, tables),
            blocks=[Block(text=line) for line in lines],
            tables=tables,
            has_text_layer=True,
            backend="docx",
        )
    ]


def _pages_from_pptx(path: Path) -> list[Page]:
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[Page] = []
    for slide in prs.slides:
        lines: list[str] = []
        tables: list[Table] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                rows = [
                    [cell.text.strip() for cell in row.cells] for row in shape.table.rows
                ]
                parsed = _table_from_rows(rows)
                if parsed:
                    tables.append(parsed)
        out.append(
            Page(
                index=0,
                image=None,
                source=path.name,
                markdown=_with_tables(lines, tables),
                blocks=[Block(text=line) for line in lines],
                tables=tables,
                has_text_layer=True,
                backend="pptx",
            )
        )
    return out


def _pages_from_xlsx(path: Path) -> list[Page]:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    out: list[Page] = []
    try:
        for sheet in wb.worksheets:
            rows = [
                ["" if c is None else str(c) for c in row]
                for row in sheet.iter_rows(values_only=True)
            ]
            table = _table_from_rows(rows)
            tables = [table] if table else []
            if table:
                table.caption = sheet.title
            out.append(
                Page(
                    index=0,
                    image=None,
                    source="%s#%s" % (path.name, sheet.title),
                    markdown=_with_tables([], tables)
                    or "\n".join("\t".join(r) for r in rows),
                    tables=tables,
                    has_text_layer=True,
                    backend="xlsx",
                )
            )
    finally:
        wb.close()
    return out


def _pages_from_text(path: Path) -> list[Page]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [
        Page(
            index=0,
            image=None,
            source=path.name,
            markdown=text,
            blocks=[Block(text=line) for line in text.splitlines() if line.strip()],
            has_text_layer=True,
            backend="text",
        )
    ]


def _expand_zip(path: Path, workdir: Path) -> list[Path]:
    """Extract a ZIP into the workdir and return the supported members.

    Members are extracted through :meth:`Path.name` only, so an archive
    containing ``../`` or absolute paths cannot escape the workdir.
    """
    target = workdir / ("unzipped_" + path.stem)
    target.mkdir(parents=True, exist_ok=True)
    out: list[Path] = []
    with zipfile.ZipFile(path) as zf:
        for info in zf.infolist():
            if info.is_dir():
                continue
            member = Path(info.filename)
            if member.suffix.lower() not in SUPPORTED_SUFFIXES:
                continue
            if member.suffix.lower() in ZIP_SUFFIXES:
                continue  # do not recurse into nested archives
            safe = target / member.name
            with zf.open(info) as src, open(safe, "wb") as dst:
                dst.write(src.read())
            out.append(safe)
    return sorted(out)


# ---------------------------------------------------------------------------
# public entry point
# ---------------------------------------------------------------------------

def _pages_for_path(path: Path, dpi: int, workdir: Path) -> list[Page]:
    suffix = path.suffix.lower()
    if suffix in ZIP_SUFFIXES:
        pages: list[Page] = []
        for member in _expand_zip(path, workdir):
            pages.extend(_pages_for_path(member, dpi, workdir))
        return pages
    if suffix in PDF_SUFFIXES:
        return _pages_from_pdf(path, dpi)
    if suffix in DOCX_SUFFIXES:
        return _pages_from_docx(path)
    if suffix in PPTX_SUFFIXES:
        return _pages_from_pptx(path)
    if suffix in XLSX_SUFFIXES:
        return _pages_from_xlsx(path)
    if suffix in TEXT_SUFFIXES:
        return _pages_from_text(path)
    if suffix in IMAGE_SUFFIXES or suffix in HEIC_SUFFIXES:
        return _pages_from_image(path)
    raise ValueError(
        "unsupported file type %r (supported: %s)"
        % (suffix or path.name, ", ".join(sorted(SUPPORTED_SUFFIXES)))
    )


def load(
    sources: str | os.PathLike | Image.Image | Sequence[Any],
    *,
    instruction: str = "",
    dpi: int = DEFAULT_DPI,
    workdir: str | os.PathLike | None = None,
) -> Document:
    """Build a Document from paths, PIL images, or a mix of both.

    Unreadable inputs are recorded as warnings rather than raising, so one bad
    file in a batch does not lose the rest of the run.
    """
    if isinstance(sources, (str, os.PathLike, Image.Image)):
        items: list[Any] = [sources]
    else:
        items = list(sources)

    work = Path(workdir) if workdir else Path(os.environ.get("DOCAGENT_WORKDIR", ".")) / ".docagent"
    work.mkdir(parents=True, exist_ok=True)

    doc = Document(instruction=instruction, workdir=str(work))
    for item in items:
        if item is None:
            continue
        if isinstance(item, Image.Image):
            doc.pages.append(Page(index=0, image=item.convert("RGB"), source="upload"))
            doc.sources.append("upload")
            continue
        # Gradio file objects expose .name / .path rather than being paths.
        raw = getattr(item, "path", None) or getattr(item, "name", None) or item
        path = Path(str(raw))
        if not path.exists():
            doc.warnings.append("missing file: %s" % path.name)
            continue
        try:
            pages = _pages_for_path(path, dpi, work)
        except Exception as exc:
            doc.warnings.append("could not read %s: %s" % (path.name, exc))
            log.warning("ingest failed for %s: %s", path, exc)
            continue
        if not pages:
            doc.warnings.append("no pages found in %s" % path.name)
            continue
        doc.pages.extend(pages)
        doc.sources.append(path.name)

    for i, page in enumerate(doc.pages):
        page.index = i
    return doc


def supported_extensions() -> list[str]:
    return sorted(SUPPORTED_SUFFIXES)
