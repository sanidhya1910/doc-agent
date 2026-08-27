"""Generate the test fixture matrix.

Run once with ``python tests/make_fixtures.py``.  Everything is synthesised
locally, so the fixtures are deterministic and the repo carries no scanned
documents of anyone real.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from PIL import Image, ImageDraw

HERE = Path(__file__).parent / "fixtures"

PROSE = (
    "The quarterly review covers the period ending 31 March 2026. Revenue rose "
    "to 4.2 million, an increase of eleven percent against the prior quarter. "
    "Operating costs were held flat despite the two new regional offices. "
    "The board approved a further hiring round of fourteen engineers. "
    "Risks remain concentrated in the supply agreement with the Kanto plant, "
    "which expires in November and has not yet been renewed."
)


def _blank(w: int = 1100, h: int = 780) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    img = Image.new("RGB", (w, h), "white")
    return img, ImageDraw.Draw(img)


def _lines(draw: ImageDraw.ImageDraw, lines: list[str], x: int = 50, y: int = 50,
           step: int = 30) -> None:
    for i, line in enumerate(lines):
        draw.text((x, y + i * step), line, fill="black")


def invoice_image() -> Image.Image:
    img, d = _blank()
    _lines(d, [
        "ACME SUPPLIES LTD", "", "Invoice number: 2026-114", "Invoice date: 2026-03-12",
        "Bill to: Northwind Trading", "", "Item            Qty        Price      Amount",
        "-------------------------------------------------", "Widget            2        9.99      19.98",
        "Gadget            1       24.50      24.50", "Bracket          10        0.40       4.00",
        "-------------------------------------------------", "Subtotal                            48.48",
        "Tax (0%)                             0.00", "Amount due:                         48.48",
    ])
    return img


def table_image() -> Image.Image:
    img, d = _blank(1000, 560)
    _lines(d, [
        "Regional headcount", "",
        "Region        2025      2026", "North           41        52",
        "South           28        27", "East            15        22", "West             9        18",
    ])
    for y in (95, 130, 165, 200, 235, 270):
        d.line((50, y, 500, y), fill="black")
    return img


def ruled_table_image() -> Image.Image:
    """A genuine ruled table: grid lines, header row, numeric cells.

    Distinct from ``table_image`` -- this one has full vertical and horizontal
    ruling, which is what the ruling-line detector keys off.
    """
    img = Image.new("RGB", (760, 330), "white")
    d = ImageDraw.Draw(img)
    cols = [40, 260, 420, 580, 720]
    rows = [40, 80, 120, 160, 200, 240]
    for x in cols:
        d.line((x, rows[0], x, rows[-1]), fill="black", width=2)
    for y in rows:
        d.line((cols[0], y, cols[-1], y), fill="black", width=2)
    data = [
        ["Item", "Qty", "Unit price", "Amount"],
        ["Widget", "2", "9.99", "19.98"],
        ["Gadget", "1", "24.50", "24.50"],
        ["Bracket", "10", "0.40", "4.00"],
        ["Spacer", "5", "1.20", "6.00"],
    ]
    for r, row in enumerate(data):
        for c, cell in enumerate(row):
            d.text((cols[c] + 10, rows[r] + 12), cell, fill="black")
    return img


def prose_image() -> Image.Image:
    img, d = _blank(1100, 500)
    words = PROSE.split()
    lines, current = [], ""
    for word in words:
        if len(current) + len(word) + 1 > 62:
            lines.append(current)
            current = word
        else:
            current = (current + " " + word).strip()
    lines.append(current)
    _lines(d, ["QUARTERLY REVIEW", ""] + lines, step=28)
    return img


def handwriting_image() -> Image.Image:
    """Stand-in for a handwritten note: irregular baselines and spacing."""
    img, d = _blank(900, 420)
    import random

    random.seed(7)
    for i, line in enumerate([
        "meeting notes - tuesday",
        "call the supplier back",
        "budget approved, 14 hires",
        "check kanto contract nov",
    ]):
        y = 60 + i * 70 + random.randint(-6, 6)
        x = 60
        for ch in line:
            d.text((x, y + random.randint(-3, 3)), ch, fill="black")
            x += 11 + random.randint(-1, 2)
    return img


def build_digital_pdf(path: Path) -> None:
    """A PDF with a real text layer -- must cost zero GPU."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, height - 60, "Quarterly Review 2026 Q1")
    c.setFont("Helvetica", 11)
    y = height - 100
    for line in PROSE.split(". "):
        if line.strip():
            c.drawString(50, y, line.strip().rstrip(".") + ".")
            y -= 18
    c.showPage()
    c.setFont("Helvetica", 11)
    c.drawString(50, height - 60, "Appendix: contact details")
    c.drawString(50, height - 85, "Owner: Operations")
    c.drawString(50, height - 105, "Reference: QR-2026-01")
    c.showPage()
    c.save()


def build_scanned_pdf(path: Path, images: list[Image.Image]) -> None:
    """Images only, no text layer -- must require OCR."""
    rgb = [im.convert("RGB") for im in images]
    rgb[0].save(str(path), save_all=True, append_images=rgb[1:], format="PDF")


def build_docx(path: Path) -> None:
    import docx

    d = docx.Document()
    d.add_heading("Supplier agreement", level=1)
    d.add_paragraph("This agreement covers the period 2026-04 through 2027-03.")
    d.add_paragraph("Renewal notice: 90 days")
    t = d.add_table(rows=1, cols=3)
    t.style = "Table Grid"
    for cell, text in zip(t.rows[0].cells, ["Term", "Value", "Unit"]):
        cell.text = text
    for row in (("Volume", "12000", "units"), ("Rate", "0.40", "each")):
        cells = t.add_row().cells
        for cell, text in zip(cells, row):
            cell.text = text
    d.save(str(path))


def build_pptx(path: Path) -> None:
    """Slides with body text and a table -- one Page per slide, no OCR."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Q1 review"
    slide.placeholders[1].text = "Revenue up eleven percent\nTwo new regional offices"

    slide = prs.slides.add_slide(blank)
    shape = slide.shapes.add_table(3, 3, Inches(0.5), Inches(0.8), Inches(6), Inches(1.8))
    rows = [["Region", "2025", "2026"], ["North", "41", "52"], ["South", "28", "27"]]
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            shape.table.cell(r, c).text = text

    prs.save(str(path))


def build_xlsx(path: Path) -> None:
    """A two-sheet workbook -- one Page per sheet, each a table."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Headcount"
    for row in (["Region", "2025", "2026"], ["North", 41, 52], ["South", 28, 27]):
        ws.append(row)

    ws2 = wb.create_sheet("Rates")
    for row in (["Term", "Value", "Unit"], ["Volume", 12000, "units"],
                ["Rate", 0.4, "each"]):
        ws2.append(row)

    wb.save(str(path))


def build_heic(path: Path, image: Image.Image) -> bool:
    """A HEIC photo, the format phone cameras actually produce.

    Returns False when the installed pillow-heif cannot encode, so the fixture
    set degrades rather than the generator failing.
    """
    try:
        import pillow_heif

        pillow_heif.register_heif_opener()
        image.convert("RGB").save(str(path), format="HEIF")
        return True
    except Exception as exc:  # noqa: BLE001
        print("  (skipped HEIC fixture: %s)" % exc)
        return False


def main() -> None:
    HERE.mkdir(parents=True, exist_ok=True)

    invoice, table, prose, hand = (
        invoice_image(), table_image(), prose_image(), handwriting_image()
    )
    ruled_table_image().save(HERE / "ruled_table.png")
    invoice.save(HERE / "invoice.png")
    table.save(HERE / "table.png")
    prose.save(HERE / "prose.png")
    hand.save(HERE / "handwriting.png")

    # Multi-frame TIFF -> one page per frame. Frames are padded to a common
    # canvas because Pillow will not encode mixed-size frames into one TIFF.
    frames = [invoice, table, prose]
    w = max(f.width for f in frames)
    h = max(f.height for f in frames)
    padded = []
    for frame in frames:
        canvas = Image.new("RGB", (w, h), "white")
        canvas.paste(frame, (0, 0))
        padded.append(canvas)
    padded[0].save(HERE / "multipage.tiff", save_all=True, append_images=padded[1:],
                   compression="tiff_deflate")

    build_digital_pdf(HERE / "digital.pdf")
    build_scanned_pdf(HERE / "scanned.pdf", [invoice, table, prose])
    build_docx(HERE / "agreement.docx")
    build_pptx(HERE / "deck.pptx")
    build_xlsx(HERE / "workbook.xlsx")
    build_heic(HERE / "photo.heic", invoice)

    with zipfile.ZipFile(HERE / "batch.zip", "w") as zf:
        for name in ("invoice.png", "table.png"):
            zf.write(HERE / name, name)

    print("fixtures written to", HERE)
    for p in sorted(HERE.iterdir()):
        print("  %-20s %8d bytes" % (p.name, p.stat().st_size))


if __name__ == "__main__":
    main()
